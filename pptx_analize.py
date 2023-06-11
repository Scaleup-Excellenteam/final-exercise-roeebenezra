import json
import aiohttp
import os
import asyncio

from openai import api_key
from pptx import Presentation
from openAi import OpenAIAPI


async def generate_ai_response(prompt):

    # Make an asynchronous request to the OpenAI API to generate AI response
    api_endpoint = "https://api.openai.com/v1/engines/davinci-codex/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    data = {
        "prompt": prompt,
        "max_tokens": 100,
        "temperature": 0.6
    }

    for _ in range(3):  # Retry 3 times
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(api_endpoint, headers=headers, json=data) as response:
                    response_data = await response.json()
            ai_reply = response_data["choices"][0]["text"].strip()
            return ai_reply
        except aiohttp.ClientError:
            await asyncio.sleep(2)  # Wait for 2 seconds before retrying

    # If all retries fail, handle the error accordingly
    raise Exception("Unable to connect to the OpenAI API")


async def process_slide(slide, slide_number):
    # Iterate through slide content
    slide_explanation = {"slide_number": slide_number, "explanation": ""}
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Extract text from shape
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    prompt = run.text.strip().replace('\xa0', ' ')

                    if prompt:
                        # Generate AI response asynchronously
                        ai_response = await generate_ai_response(prompt)

                        # Append AI response to slide explanation
                        slide_explanation["explanation"] += ai_response + " "

    return slide_explanation


async def parse_pptx_file(filepath):
    presentation = Presentation(filepath)
    slide_tasks = []

    # Create a list of tasks for processing each slide asynchronously
    for i, slide in enumerate(presentation.slides):
        slide_number = i + 1
        slide_task = asyncio.create_task(process_slide(slide, slide_number))
        slide_tasks.append(slide_task)

    # Wait for all tasks to complete
    slide_explanations = await asyncio.gather(*slide_tasks)

    # Save slide explanations to a JSON file with the same name as the original presentation
    file_name = os.path.splitext(os.path.basename(filepath))[0]
    output_file = f"{file_name}.json"
    with open(output_file, "w") as json_file:
        json.dump(slide_explanations, json_file)

    print(f"Slide explanations saved to {output_file}")


if __name__ == "__main__":
    # Path to .pptx file
    pptx_file_path = "/Users/roeebenezra/PycharmProjects/final-exercise-roeebenezra/test.pptx"

    # Run the asynchronous parsing
    asyncio.run(parse_pptx_file(pptx_file_path))
