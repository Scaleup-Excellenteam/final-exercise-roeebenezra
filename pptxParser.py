from pptx import Presentation


class PresentationParser:
    """
    Class to parse a PowerPoint presentation.
    """
    def __init__(self, presentation_path):
        self.presentation_path = presentation_path

    def process_presentation(self) -> list[str]:
        """
        Process each slide in the presentation and retrieve the text content.
        :return: list[str]: List of text content for each slide.
        """
        prs = Presentation(self.presentation_path)
        slides = [self.process_slide_text(slide) for slide in prs.slides]
        return slides

    @staticmethod
    def process_slide_text(slide) -> str:
        """
        Extract the text content from a slide in the PowerPoint presentation.
        :param slide: Slide object from the PowerPoint presentation.
        :return: str: Processed text content from the slide.
        """
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text.strip())
        return " ".join(slide_text)
